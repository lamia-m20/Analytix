from io import BytesIO
import csv
import io
import os
import zipfile
from unittest.mock import MagicMock, patch

import pandas as pd
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches
from PIL import Image as PILImage
from django.contrib.auth import get_user_model
from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import TestCase
from django.urls import reverse

from dashboards.models import Dashboard
from dashboards.models import DashboardWidget
from dashboards.services.dashboard_ai import (
    DashboardAIUnavailableError,
)
from dashboards.services.dashboard_ai import (
    DashboardPlanValidationError,
)
from dashboards.services.dashboard_ai import (
    apply_dashboard_plan,
)
from dashboards.services.dashboard_ai import (
    request_dashboard_plan,
)
from dashboards.services.dashboard_store import (
    get_or_create_dataset_dashboard,
)
from analysis.models import AnalysisJob, AnalysisResult

from .models import Dataset
from .models import DatasetColumn
from .models import DatasetSheet


class DatasetExportTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(username='exporter', password='test-password')
        self.other = get_user_model().objects.create_user(username='other-exporter', password='test-password')
        self.dataset = Dataset.objects.create(
            user=self.user, title='تقرير المبيعات', original_filename='المبيعات.xlsx',
            file='tests/sales.xlsx', status='ready',
        )
        self.frames = {
            'المبيعات': pd.DataFrame({'المدينة': ['الرياض', 'جدة'], 'القيمة': [100, 200]}),
            'الفروع': pd.DataFrame({'الفرع': ['شمال'], 'الموظفون': [8]}),
        }
        self.client.force_login(self.user)

    def _response(self, name):
        with patch('datasets.views.load_workbook_data', return_value=self.frames):
            return self.client.get(reverse(f'datasets:{name}', args=[self.dataset.pk]))

    def test_anonymous_user_is_redirected(self):
        self.client.logout()
        response = self.client.get(reverse('datasets:export_excel', args=[self.dataset.pk]))
        self.assertEqual(response.status_code, 302)
        self.assertIn('/accounts/login/', response.url)

    def test_user_cannot_export_another_users_dataset(self):
        self.client.force_login(self.other)
        response = self.client.get(reverse('datasets:export_pdf', args=[self.dataset.pk]))
        self.assertEqual(response.status_code, 404)

    def test_anonymous_user_cannot_export_powerpoint(self):
        self.client.logout()
        response = self.client.get(
            reverse('datasets:export_powerpoint', args=[self.dataset.pk])
        )
        self.assertEqual(response.status_code, 302)
        self.assertIn('/accounts/login/', response.url)

    def test_user_cannot_export_another_users_powerpoint(self):
        self.client.force_login(self.other)
        response = self.client.get(
            reverse('datasets:export_powerpoint', args=[self.dataset.pk])
        )
        self.assertEqual(response.status_code, 404)

    def test_missing_dataset_returns_404(self):
        response = self.client.get(reverse('datasets:export_csv', args=[999999]))
        self.assertEqual(response.status_code, 404)

    def test_excel_export(self):
        response = self._response('export_excel')
        content = b''.join(response.streaming_content)
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response['Content-Type'], 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        self.assertIn('attachment;', response['Content-Disposition'])
        self.assertIn('analytix_report_', response['Content-Disposition'])
        self.assertTrue(content.startswith(b'PK'))

    def test_csv_export_uses_zip_for_multiple_sheets_and_arabic_utf8(self):
        response = self._response('export_csv')
        content = b''.join(response.streaming_content)
        self.assertEqual(response['Content-Type'], 'application/zip')
        self.assertIn('.zip', response['Content-Disposition'])
        with zipfile.ZipFile(BytesIO(content)) as archive:
            csv_content = archive.read(archive.namelist()[0])
        self.assertTrue(csv_content.startswith(b'\xef\xbb\xbf'))
        self.assertIn('الرياض'.encode('utf-8'), csv_content)
        header = csv_content.decode('utf-8-sig').splitlines()[0]
        self.assertIn(';', header)
        self.assertNotIn(',', header)

    def test_single_sheet_csv_uses_excel_arabic_dialect(self):
        frames = {
            'المبيعات': pd.DataFrame({
                'المدينة': ['الرياض', 'جدة'],
                'القيمة': [100, 200],
            })
        }
        with patch('datasets.views.load_workbook_data', return_value=frames):
            response = self.client.get(
                reverse('datasets:export_csv', args=[self.dataset.pk])
            )
        content = b''.join(response.streaming_content)
        self.assertEqual(response['Content-Type'], 'text/csv; charset=utf-8-sig')
        self.assertTrue(content.startswith(b'\xef\xbb\xbf'))
        decoded = content.decode('utf-8-sig')
        parsed = list(csv.reader(io.StringIO(decoded, newline=''), delimiter=';'))
        self.assertEqual(parsed[0], ['المدينة', 'القيمة'])
        self.assertEqual(len(parsed[1]), 2)

    def test_pdf_export(self):
        from datasets.exporters import get_arabic_font_path
        try:
            get_arabic_font_path()
        except Exception:
            self.skipTest('ضع خطًا عربيًا في static/fonts لتشغيل اختبار PDF الفعلي.')
        response = self._response('export_pdf')
        content = b''.join(response.streaming_content)
        self.assertEqual(response['Content-Type'], 'application/pdf')
        self.assertIn('attachment;', response['Content-Disposition'])
        self.assertTrue(content.startswith(b'%PDF'))
        self.assertNotIn(b'nnnnnn', content.lower())
        self.assertIn(b'/Subtype /Image', content)

    def test_powerpoint_export_is_valid_and_contains_core_slides(self):
        response = self._response('export_powerpoint')
        content = b''.join(response.streaming_content)
        self.assertEqual(
            response['Content-Type'],
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        )
        self.assertIn('attachment;', response['Content-Disposition'])
        self.assertIn('.pptx', response['Content-Disposition'])
        self.assertTrue(content.startswith(b'PK'))
        presentation = Presentation(BytesIO(content))
        self.assertGreaterEqual(len(presentation.slides), 8)
        all_text = '\n'.join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, 'text')
        )
        self.assertIn('تقرير تحليل البيانات الشامل', all_text)
        self.assertIn('الملخص التنفيذي', all_text)
        self.assertIn('التوصيات', all_text)
        self.assertIn('مؤشرات الجودة', all_text)
        self.assertIn('Analytix', all_text)
        picture_count = sum(
            1 for slide in presentation.slides for shape in slide.shapes
            if shape.shape_type == 13
        )
        self.assertGreaterEqual(picture_count, 1)

    def test_powerpoint_without_numeric_columns_or_widgets(self):
        frames = {'نصوص': pd.DataFrame({'الفئة': ['أ', 'ب', 'أ']})}
        with patch('datasets.views.load_workbook_data', return_value=frames):
            response = self.client.get(
                reverse('datasets:export_powerpoint', args=[self.dataset.pk])
            )
        content = b''.join(response.streaming_content)
        self.assertTrue(content.startswith(b'PK'))
        self.assertGreater(len(Presentation(BytesIO(content)).slides), 5)

    def test_powerpoint_chart_failure_does_not_stop_export(self):
        with patch('datasets.views.load_workbook_data', return_value=self.frames), patch(
            'datasets.powerpoint._dashboard_chart', return_value=None
        ):
            response = self.client.get(
                reverse('datasets:export_powerpoint', args=[self.dataset.pk])
            )
        content = b''.join(response.streaming_content)
        self.assertEqual(response.status_code, 200)
        self.assertTrue(content.startswith(b'PK'))
        self.assertGreater(len(Presentation(BytesIO(content)).slides), 5)

    def test_powerpoint_uses_design_grid_and_six_summary_cards(self):
        response = self._response('export_powerpoint')
        presentation = Presentation(BytesIO(b''.join(response.streaming_content)))
        self.assertAlmostEqual(presentation.slide_width / presentation.slide_height, 16 / 9, places=2)
        summary_slide = presentation.slides[1]
        kpi_cards = [shape for shape in summary_slide.shapes if shape.name == 'KPI_CARD']
        self.assertEqual(len(kpi_cards), 6)
        for slide in list(presentation.slides)[1:-1]:
            names = {shape.name for shape in slide.shapes}
            self.assertIn('ANALYTIX_HEADER', names)
            self.assertIn('ANALYTIX_FOOTER', names)
            self.assertIn('PAGE_NUMBER', names)
        chart_pictures = [
            shape for slide in presentation.slides for shape in slide.shapes
            if shape.name == 'CHART_PICTURE'
        ]
        self.assertTrue(chart_pictures)
        self.assertGreaterEqual(chart_pictures[0].width, Inches(8.5))
        self.assertGreaterEqual(chart_pictures[0].height, Inches(3.8))
        chart_slide_text = '\n'.join(
            shape.text
            for slide in presentation.slides
            if any(shape.name == 'CHART_PICTURE' for shape in slide.shapes)
            for shape in slide.shapes
            if hasattr(shape, 'text')
        )
        self.assertIn('القراءة الرئيسية', chart_slide_text)
        for slide in presentation.slides:
            if any(shape.name == 'CHART_PICTURE' for shape in slide.shapes):
                insight_panels = [shape for shape in slide.shapes if shape.name == 'INSIGHT_PANEL']
                self.assertEqual(len(insight_panels), 1)
                self.assertEqual(len([shape for shape in slide.shapes if shape.name == 'CHART_KPI']), 4)
                background = slide.background.fill.fore_color.rgb
                self.assertEqual(str(background), '08121E')
        executive_names = {shape.name for shape in summary_slide.shapes}
        self.assertIn('DASHBOARD_SIDEBAR', executive_names)
        self.assertIn('EXECUTIVE_BAR', executive_names)
        self.assertIn('EXECUTIVE_DONUT', executive_names)
        self.assertIn('EXECUTIVE_GAUGE', executive_names)
        dashboard_slides = list(presentation.slides)[1:-1]
        self.assertTrue(all(
            any(shape.name == 'DASHBOARD_SIDEBAR' for shape in slide.shapes)
            for slide in dashboard_slides
        ))
        self.assertTrue(any(
            shape.name in {'QUALITY_BAR', 'COMPLETION_DONUT', 'QUALITY_GAUGE'}
            for slide in presentation.slides for shape in slide.shapes
        ))
        self.assertTrue(any(
            shape.name == 'SHEET_DASHBOARD_CHART'
            for slide in presentation.slides for shape in slide.shapes
        ))
        for slide in presentation.slides:
            for shape in slide.shapes:
                self.assertGreaterEqual(shape.left, 0)
                self.assertGreaterEqual(shape.top, 0)
                self.assertLessEqual(shape.left + shape.width, presentation.slide_width)
                self.assertLessEqual(shape.top + shape.height, presentation.slide_height)

    def test_powerpoint_recommendations_use_one_priority_map(self):
        from datasets.exporters import build_report_data
        from datasets.powerpoint import build_powerpoint_report
        report = build_report_data(self.dataset, self.frames)
        report['executive_summary']['recommendations'] = [
            f'توصية عملية طويلة رقم {index} مرتبطة بنتائج جودة البيانات.'
            for index in range(7)
        ]
        presentation = Presentation(build_powerpoint_report(self.dataset, report))
        recommendation_slides = []
        for slide in presentation.slides:
            text = '\n'.join(shape.text for shape in slide.shapes if hasattr(shape, 'text'))
            if 'خارطة الأولويات' in text:
                recommendation_slides.append(slide)
        self.assertEqual(len(recommendation_slides), 1)
        priority_text = '\n'.join(
            shape.text for shape in recommendation_slides[0].shapes if hasattr(shape, 'text')
        )
        self.assertIn('مرتفعة', priority_text)
        self.assertIn('متوسطة', priority_text)
        self.assertIn('منخفضة', priority_text)

    def test_powerpoint_chart_arabic_is_shaped_only_for_matplotlib(self):
        from datasets.powerpoint import prepare_arabic_for_chart

        samples = ('درجة الجودة', 'القيم الفارغة', 'الصفوف المكررة', 'جودة البيانات 79%')
        for sample in samples:
            prepared = prepare_arabic_for_chart(sample)
            self.assertNotEqual(prepared, sample)
            self.assertNotIn('nnnn', prepared.lower())
        self.assertIn('79%', prepare_arabic_for_chart('جودة البيانات 79%'))
        self.assertEqual(prepare_arabic_for_chart('Analytix'), 'Analytix')

    def test_pdf_missing_arabic_font_returns_clear_message(self):
        from datasets.exporters import ArabicFontNotFoundError
        with patch('datasets.exporters.get_arabic_font_path', side_effect=ArabicFontNotFoundError('الخط العربي غير موجود')):
            response = self._response('export_pdf')
        self.assertEqual(response.status_code, 503)
        self.assertContains(response, 'الخط العربي غير موجود', status_code=503)

    def test_prepare_arabic_text_preserves_english(self):
        from datasets.exporters import prepare_arabic_text
        samples = (
            'تقرير تحليل البيانات', 'القيم الفارغة', 'عدد الصفوف',
            'عدد الأعمدة', 'Analytix', 'تقرير Analytix لتحليل البيانات',
        )
        prepared = [prepare_arabic_text(value) for value in samples]
        self.assertEqual(prepared[4], 'Analytix')
        self.assertIn('Analytix', prepared[5])
        self.assertNotIn('nnnnnn', ''.join(prepared).lower())

    def test_cloud_file_failure_returns_arabic_message(self):
        from datasets.exporters import ExportSourceError
        with patch('datasets.views.load_workbook_data', side_effect=ExportSourceError('تعذر الوصول إلى ملف البيانات.')):
            response = self.client.get(reverse('datasets:export_excel', args=[self.dataset.pk]))
        self.assertEqual(response.status_code, 503)
        self.assertContains(response, 'تعذر الوصول', status_code=503)

    def test_unified_report_data_covers_quality_duplicates_dates_and_arabic(self):
        from datasets.exporters import build_report_data
        frames = {
            'بيانات عربية': pd.DataFrame({
                'المدينة': ['الرياض', 'الرياض', None, 'جدة'],
                'القيمة': [10, 10, 1000, 20],
                'التاريخ': pd.to_datetime(['2025-01-01', '2025-01-01', '2025-03-01', '2025-04-01']),
            })
        }
        report = build_report_data(self.dataset, frames)
        sheet = report['sheets'][0]
        self.assertEqual(report['summary']['sheet_count'], 1)
        self.assertEqual(sheet['duplicates'], 1)
        self.assertEqual(sheet['missing'][0]['column'], 'المدينة')
        self.assertEqual(sheet['date_analysis'][0]['column'], 'التاريخ')
        self.assertLess(sheet['quality_score'], 100)
        self.assertTrue(report['executive_summary']['insights'])

    def test_report_without_numeric_columns(self):
        from datasets.exporters import build_report_data
        report = build_report_data(self.dataset, {
            'نصوص': pd.DataFrame({'الفئة': ['أ', 'ب', 'أ']})
        })
        self.assertEqual(report['sheets'][0]['numeric_stats'], [])
        self.assertTrue(report['sheets'][0]['text_analysis'])

    def test_comprehensive_excel_contains_required_sheets(self):
        from datasets.exporters import build_excel, build_report_data
        report = build_report_data(self.dataset, self.frames)
        stream = build_excel(self.dataset, report)
        workbook = load_workbook(stream, read_only=True)
        required = {
            'غلاف التقرير', 'Dashboard', 'الملخص التنفيذي', 'مؤشرات الجودة', 'معلومات الأوراق',
            'جودة البيانات', 'القيم الفارغة', 'الصفوف المكررة', 'أنواع الأعمدة',
            'الإحصاءات الرقمية', 'التحليل النصي', 'تحليل التواريخ',
            'القيم الشاذة', 'الارتباطات', 'عناصر الداشبورد', 'التوصيات',
        }
        self.assertTrue(required.issubset(set(workbook.sheetnames)))
        workbook.close(); stream.close()

    def test_shared_chart_theme_uses_high_resolution(self):
        from datasets.exporters import _report_chart
        chart = _report_chart(['أ', 'ب', 'ج'], [10, 20, 15], 'مؤشر الجودة')
        image = PILImage.open(chart)
        dpi = image.info.get('dpi', (0, 0))
        self.assertGreaterEqual(dpi[0], 299)
        self.assertGreater(image.width, 2000)
        image.close(); chart.close()

    def test_pdf_chart_with_many_categories_stays_vertical(self):
        from datasets.exporters import _report_chart

        figure = MagicMock()
        axis = MagicMock()
        axis.get_xticklabels.return_value = []
        axis.get_yticklabels.return_value = []
        with patch('datasets.exporters.plt.subplots', return_value=(figure, axis)):
            chart = _report_chart(
                ['فئة 1', 'فئة 2', 'فئة 3', 'فئة 4', 'فئة 5', 'فئة 6'],
                [1, 2, 3, 4, 5, 6],
                'المقارنة',
            )

        axis.bar.assert_called_once()
        axis.barh.assert_not_called()
        axis.set_xticklabels.assert_called_once()
        chart.close()

    def test_powerpoint_chart_with_many_categories_stays_vertical(self):
        from datasets.powerpoint import _dashboard_chart

        figure = MagicMock()
        axis = MagicMock()
        axis.get_xticklabels.return_value = []
        axis.get_yticklabels.return_value = []
        axis.bar_label.return_value = []
        with patch('datasets.powerpoint.plt.subplots', return_value=(figure, axis)):
            chart = _dashboard_chart(
                ['فئة 1', 'فئة 2', 'فئة 3', 'فئة 4', 'فئة 5', 'فئة 6'],
                [1, 2, 3, 4, 5, 6],
            )

        axis.bar.assert_called_once()
        axis.barh.assert_not_called()
        axis.set_xticklabels.assert_called_once()
        chart.close()

    def test_zip_contains_analysis_files_and_readme(self):
        from datasets.exporters import build_csv, build_report_data
        stream, content_type, extension = build_csv(
            self.dataset, build_report_data(self.dataset, self.frames)
        )
        with zipfile.ZipFile(stream) as archive:
            names = set(archive.namelist())
        self.assertEqual(content_type, 'application/zip')
        self.assertEqual(extension, '.zip')
        self.assertIn('README.txt', names)
        self.assertIn('الملخص_التنفيذي.csv', names)
        self.assertIn('جودة_البيانات.csv', names)
        with zipfile.ZipFile(stream) as archive:
            for name in (item for item in archive.namelist() if item.endswith('.csv')):
                csv_bytes = archive.read(name)
                self.assertTrue(csv_bytes.startswith(b'\xef\xbb\xbf'))
                first_line = csv_bytes.decode('utf-8-sig').splitlines()[0]
                if first_line:
                    self.assertNotIn(',', first_line)

    def test_export_formats_have_distinct_roles_and_payloads(self):
        from datasets.exporters import (
            build_csv_package, build_excel_report, build_pdf_report,
            build_report_data,
        )
        one_sheet = {'المبيعات': self.frames['المبيعات']}
        report = build_report_data(self.dataset, one_sheet)
        excel = build_excel_report(self.dataset, report).getvalue()
        pdf = build_pdf_report(self.dataset, report).getvalue()
        csv_stream, _, _ = build_csv_package(self.dataset, report)
        csv_payload = csv_stream.getvalue()
        self.assertTrue(excel.startswith(b'PK'))
        self.assertTrue(pdf.startswith(b'%PDF'))
        self.assertTrue(csv_payload.startswith(b'\xef\xbb\xbf'))
        self.assertNotEqual(excel, pdf)
        self.assertNotEqual(pdf, csv_payload)

    def test_chart_failure_does_not_stop_pdf(self):
        from datasets.exporters import build_pdf, build_report_data
        report = build_report_data(self.dataset, self.frames)
        with patch('datasets.exporters.plt.subplots', side_effect=RuntimeError('chart failed')):
            stream = build_pdf(self.dataset, report)
        self.assertTrue(stream.read(4).startswith(b'%PDF'))
        stream.close()

    def test_large_report_limit_is_documented(self):
        from datasets.exporters import MAX_REPORT_CELLS, PREVIEW_ROWS
        self.assertEqual(PREVIEW_ROWS, 100)
        self.assertGreaterEqual(MAX_REPORT_CELLS, 1_000_000)

    def test_unnamed_columns_are_cleaned_for_report_only(self):
        from datasets.exporters import build_report_data
        source = pd.DataFrame({
            'الاسم': ['أ', 'ب'],
            'Unnamed: 1': [None, None],
            'Unnamed: 2': ['قيمة', None],
        })
        report = build_report_data(self.dataset, {'ورقة': source})
        names = [item['column'] for item in report['sheets'][0]['columns_info']]
        self.assertNotIn('Unnamed: 1', names)
        self.assertNotIn('Unnamed: 2', names)
        self.assertIn('عمود غير مسمى 1', names)
        self.assertIn('Unnamed: 1', source.columns)
        self.assertIn('Unnamed: 2', source.columns)

    def test_excel_auto_width_dates_long_arabic_and_safe_sheet_names(self):
        from datasets.exporters import build_excel_report, build_report_data
        long_text = 'هذا نص عربي طويل لاختبار التفاف النص وضبط ارتفاع الصف تلقائيًا ' * 3
        common_prefix = 'اسم ورقة عربي طويل جدًا يتجاوز واحدًا وثلاثين حرفًا'
        frames = {
            f'{common_prefix} أ': pd.DataFrame({
                'التاريخ الطويل': pd.to_datetime(['2026-08-05']),
                'الوصف': [long_text],
                'Unnamed: 1': [None],
            }),
            f'{common_prefix} ب': pd.DataFrame({f'عمود {index}': [index] for index in range(20)}),
        }
        report = build_report_data(self.dataset, frames)
        stream = build_excel_report(self.dataset, report)
        workbook = load_workbook(stream)
        preview_names = [name for name in workbook.sheetnames if name.startswith('معاينة')]
        self.assertEqual(len(preview_names), 2)
        self.assertEqual(len(set(preview_names)), 2)
        self.assertTrue(all(len(name) <= 31 for name in preview_names))
        first_preview = workbook[preview_names[0]]
        headers = [cell.value for cell in first_preview[1]]
        self.assertNotIn('Unnamed: 1', headers)
        date_column = headers.index('التاريخ الطويل') + 1
        self.assertEqual(first_preview.cell(2, date_column).number_format, 'yyyy-mm-dd')
        self.assertGreaterEqual(first_preview.column_dimensions[get_column_letter(date_column)].width, 12)
        self.assertLessEqual(first_preview.column_dimensions['B'].width, 45)
        self.assertGreater(first_preview.row_dimensions[2].height, 15)
        self.assertTrue(first_preview.freeze_panes == 'A2')
        self.assertTrue(first_preview.auto_filter.ref)
        workbook.close(); stream.close()

    def test_preview_excludes_fully_empty_rows_and_columns(self):
        from datasets.exporters import build_report_data
        frame = pd.DataFrame({'بيانات': ['قيمة', None], 'فارغ': [None, None]})
        report = build_report_data(self.dataset, {'ورقة': frame})
        preview = report['sheets'][0]['preview'].dropna(axis=0, how='all').dropna(axis=1, how='all')
        self.assertEqual(preview.shape, (1, 1))


class DatasetDashboardTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username='analyst',
            password='strong-test-password',
        )
        self.client.force_login(self.user)

    def _excel_file(self):
        buffer = BytesIO()

        with pd.ExcelWriter(
            buffer,
            engine='openpyxl',
        ) as writer:
            pd.DataFrame(
                {
                    'المبيعات': [100, 200, None],
                    'المنطقة': ['الرياض', 'جدة', 'الرياض'],
                }
            ).to_excel(
                writer,
                sheet_name='المبيعات',
                index=False,
            )
            pd.DataFrame(
                {
                    'الموظفون': [8, 12],
                }
            ).to_excel(
                writer,
                sheet_name='الفروع',
                index=False,
            )

        return SimpleUploadedFile(
            'بيانات.xlsx',
            buffer.getvalue(),
            content_type=(
                'application/vnd.openxmlformats-officedocument.'
                'spreadsheetml.sheet'
            ),
        )

    def _dataset_with_structure(self, user=None):
        dataset = Dataset.objects.create(
            user=user or self.user,
            title='ملف المبيعات',
            file='test/sales.xlsx',
            status='uploaded',
        )
        sheet = DatasetSheet.objects.create(
            dataset=dataset,
            name='المبيعات',
            row_count=3,
            column_count=2,
        )
        DatasetColumn.objects.bulk_create([
            DatasetColumn(
                sheet=sheet,
                name='الشهر',
                position=0,
                data_type='object',
            ),
            DatasetColumn(
                sheet=sheet,
                name='القيمة',
                position=1,
                data_type='float64',
            ),
        ])
        return dataset

    def test_dashboard_page_loads(self):
        response = self.client.get(
            reverse('datasets:list')
        )

        self.assertEqual(response.status_code, 200)
        self.assertContains(response, 'ابدأ التحليل')

    def test_excel_upload_creates_analysis_context(self):
        response = self.client.post(
            reverse('datasets:list'),
            {'excel_file': self._excel_file()},
        )

        self.assertEqual(response.status_code, 200)
        self.assertTrue(response.context['analysis_complete'])
        self.assertEqual(response.context['sheets_count'], 2)
        self.assertEqual(response.context['total_rows'], 5)
        self.assertEqual(response.context['total_columns'], 3)
        self.assertEqual(
            response.context['total_missing_values'],
            1,
        )
        self.assertContains(response, 'لوحة المعلومات')

    def _saved_analysis(self, dataset):
        sheet = dataset.sheets.first()
        job = AnalysisJob.objects.create(
            owner=dataset.user,
            dataset=dataset,
            sheet=sheet,
            name=dataset.title,
            status='completed',
            progress=100,
        )
        return AnalysisResult.objects.create(
            analysis_job=job,
            summary={
                'analysis_complete': True,
                'file_name': 'sales.xlsx',
                'sheets_count': 1,
                'total_rows': 3,
                'total_columns': 2,
                'total_missing_values': 0,
                'total_numeric_columns': 1,
                'sheets_analysis': [],
                'chart_sheet_names': [],
                'chart_rows': [],
                'chart_columns': [],
                'chart_missing_values': [],
                'numeric_chart_labels': [],
                'numeric_chart_means': [],
                'numeric_chart_max_values': [],
                'numeric_chart_min_values': [],
            },
        )

    def test_my_analyses_requires_login_and_only_lists_owner_data(self):
        own = self._dataset_with_structure()
        self._saved_analysis(own)
        other = get_user_model().objects.create_user(username='other')
        foreign = self._dataset_with_structure(user=other)
        self._saved_analysis(foreign)

        response = self.client.get(reverse('datasets:my_analyses'))
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, 'sales.xlsx')
        listed_ids = [item['dataset'].pk for item in response.context['items']]
        self.assertEqual(listed_ids, [own.pk])

        self.client.logout()
        response = self.client.get(reverse('datasets:my_analyses'))
        self.assertEqual(response.status_code, 302)

    def test_saved_analysis_opens_without_rereading_excel_or_new_records(self):
        dataset = self._dataset_with_structure()
        self._saved_analysis(dataset)
        get_or_create_dataset_dashboard(dataset, has_numeric_columns=True)
        dataset_count = Dataset.objects.count()
        widget_count = DashboardWidget.objects.count()

        with patch('datasets.views._analyze_excel') as analyze:
            response = self.client.get(
                reverse('datasets:detail', args=[dataset.pk])
            )
            refresh = self.client.get(
                reverse('datasets:detail', args=[dataset.pk])
            )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(refresh.status_code, 200)
        analyze.assert_not_called()
        self.assertEqual(Dataset.objects.count(), dataset_count)
        self.assertEqual(DashboardWidget.objects.count(), widget_count)

    def test_persisting_same_upload_is_idempotent(self):
        from datasets.views import _save_analysis_result

        dataset = self._dataset_with_structure()
        context = self._saved_analysis(dataset).summary
        AnalysisJob.objects.filter(dataset=dataset).delete()

        _save_analysis_result(dataset, context)
        first_widget_count = DashboardWidget.objects.count()
        _save_analysis_result(dataset, context)

        self.assertEqual(AnalysisJob.objects.filter(dataset=dataset).count(), 1)
        self.assertEqual(
            AnalysisResult.objects.filter(analysis_job__dataset=dataset).count(),
            1,
        )
        self.assertEqual(DashboardWidget.objects.count(), first_widget_count)

    def test_other_user_cannot_open_saved_analysis(self):
        dataset = self._dataset_with_structure()
        self._saved_analysis(dataset)
        other = get_user_model().objects.create_user(username='intruder')
        self.client.force_login(other)
        response = self.client.get(
            reverse('datasets:detail', args=[dataset.pk])
        )
        self.assertEqual(response.status_code, 404)

    def test_non_excel_file_is_rejected(self):
        response = self.client.post(
            reverse('datasets:list'),
            {
                'excel_file': SimpleUploadedFile(
                    'data.txt',
                    b'not an excel file',
                ),
            },
        )

        self.assertContains(
            response,
            'صيغة الملف غير مدعومة',
        )

    def test_dashboard_edit_request_displays_confirmation(self):
        dataset = self._dataset_with_structure()
        with patch(
            'datasets.views.request_dashboard_plan',
            return_value={
                'actions': [{
                    'action': 'add',
                    'title': 'المبيعات حسب الشهر',
                    'widget_type': 'chart',
                    'chart_type': 'bar',
                    'sheet_name': 'المبيعات',
                    'x_column': 'الشهر',
                    'y_column': 'القيمة',
                    'aggregation': 'sum',
                }],
            },
        ):
            response = self.client.post(
                reverse(
                    'datasets:edit_dashboard',
                    args=[dataset.pk],
                ),
                {
                    'dashboard_request': (
                        'أضف مخططًا للمبيعات حسب الشهر'
                    ),
                },
                follow=True,
            )

        self.assertRedirects(
            response,
            reverse(
                'datasets:detail',
                args=[dataset.pk],
            ),
        )
        self.assertContains(
            response,
            'تم تعديل 1 مخططات بنجاح.',
        )

    def test_dashboard_edit_endpoint_rejects_get(self):
        dataset = Dataset.objects.create(
            user=self.user,
            title='ملف المبيعات',
            file='test/sales.xlsx',
            status='uploaded',
        )
        response = self.client.get(
            reverse(
                'datasets:edit_dashboard',
                args=[dataset.pk],
            )
        )

        self.assertEqual(response.status_code, 405)

    def test_user_cannot_edit_another_users_dataset(self):
        other_user = get_user_model().objects.create_user(
            username='other-analyst',
            password='strong-test-password',
        )
        dataset = Dataset.objects.create(
            user=other_user,
            title='ملف خاص',
            file='test/private.xlsx',
            status='uploaded',
        )

        response = self.client.post(
            reverse(
                'datasets:edit_dashboard',
                args=[dataset.pk],
            ),
            {'dashboard_request': 'احذف المخطط'},
        )

        self.assertEqual(response.status_code, 404)

    def test_empty_dashboard_edit_request_is_rejected(self):
        dataset = Dataset.objects.create(
            user=self.user,
            title='ملف المبيعات',
            file='test/sales.xlsx',
            status='uploaded',
        )

        response = self.client.post(
            reverse(
                'datasets:edit_dashboard',
                args=[dataset.pk],
            ),
            {'dashboard_request': '   '},
            follow=True,
        )

        self.assertRedirects(
            response,
            reverse(
                'datasets:detail',
                args=[dataset.pk],
            ),
        )
        self.assertContains(
            response,
            'يرجى كتابة طلب تعديل الداشبورد.',
        )

    def test_dashboard_and_default_widgets_are_created_once(self):
        dataset = Dataset.objects.create(
            user=self.user,
            title='ملف المبيعات',
            file='test/sales.xlsx',
            status='ready',
        )

        first_dashboard = get_or_create_dataset_dashboard(
            dataset,
            has_numeric_columns=True,
        )
        first_widget_ids = list(
            first_dashboard.widgets.values_list(
                'pk',
                flat=True,
            )
        )

        second_dashboard = get_or_create_dataset_dashboard(
            dataset,
            has_numeric_columns=True,
        )
        second_widget_ids = list(
            second_dashboard.widgets.values_list(
                'pk',
                flat=True,
            )
        )

        self.assertEqual(first_dashboard.pk, second_dashboard.pk)
        self.assertEqual(first_widget_ids, second_widget_ids)
        self.assertEqual(
            Dashboard.objects.filter(
                owner=self.user,
                layout_settings__dataset_id=dataset.pk,
            ).count(),
            1,
        )
        self.assertEqual(len(first_widget_ids), 4)

    def test_saved_widget_changes_survive_dashboard_reload(self):
        dataset = Dataset.objects.create(
            user=self.user,
            title='ملف المبيعات',
            file='test/sales.xlsx',
            status='ready',
        )
        dashboard = get_or_create_dataset_dashboard(
            dataset,
            has_numeric_columns=True,
        )
        widget = dashboard.widgets.get(
            settings__source='missing_values'
        )
        widget.widget_type = 'bar'
        widget.save(update_fields=['widget_type'])

        reloaded_dashboard = get_or_create_dataset_dashboard(
            dataset,
            has_numeric_columns=True,
        )
        reloaded_widget = DashboardWidget.objects.get(
            dashboard=reloaded_dashboard,
            settings__source='missing_values',
        )

        self.assertEqual(reloaded_widget.widget_type, 'bar')
        self.assertEqual(reloaded_dashboard.widgets.count(), 4)

    def test_user_cannot_open_another_users_dataset_dashboard(self):
        other_user = get_user_model().objects.create_user(
            username='dashboard-owner',
            password='strong-test-password',
        )
        dataset = Dataset.objects.create(
            user=other_user,
            title='ملف خاص',
            file='test/private.xlsx',
            status='ready',
        )

        response = self.client.get(
            reverse(
                'datasets:detail',
                args=[dataset.pk],
            )
        )

        self.assertEqual(response.status_code, 404)

    def test_anonymous_user_cannot_edit_dashboard(self):
        dataset = self._dataset_with_structure()
        self.client.logout()

        response = self.client.post(
            reverse(
                'datasets:edit_dashboard',
                args=[dataset.pk],
            ),
            {'dashboard_request': 'احذف المخطط'},
        )

        self.assertEqual(response.status_code, 302)
        self.assertIn('/accounts/login/', response.url)

    def test_openai_failure_keeps_existing_widgets(self):
        dataset = self._dataset_with_structure()
        dashboard = get_or_create_dataset_dashboard(
            dataset,
            has_numeric_columns=True,
        )
        widget_ids = list(
            dashboard.widgets.values_list(
                'pk',
                flat=True,
            )
        )

        with patch(
            'datasets.views.request_dashboard_plan',
            side_effect=DashboardAIUnavailableError(
                'تعذر الاتصال بخدمة تعديل الداشبورد حاليًا.'
            ),
        ):
            response = self.client.post(
                reverse(
                    'datasets:edit_dashboard',
                    args=[dataset.pk],
                ),
                {'dashboard_request': 'غيّر المخطط'},
                follow=True,
            )

        self.assertContains(
            response,
            'تعذر الاتصال بخدمة تعديل الداشبورد حاليًا.',
        )
        self.assertEqual(
            list(
                dashboard.widgets.values_list(
                    'pk',
                    flat=True,
                )
            ),
            widget_ids,
        )


class DashboardAIServiceTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username='ai-analyst',
            password='strong-test-password',
        )
        self.dataset = Dataset.objects.create(
            user=self.user,
            title='المبيعات',
            file='test/sales.xlsx',
            status='ready',
        )
        self.dashboard = get_or_create_dataset_dashboard(
            self.dataset,
            has_numeric_columns=True,
        )
        self.metadata = {
            'sheets': [{
                'name': 'المبيعات',
                'columns': [
                    {'name': 'الشهر', 'type': 'object'},
                    {'name': 'القيمة', 'type': 'float64'},
                ],
            }],
            'widgets': [],
        }

    def _apply(self, actions):
        return apply_dashboard_plan(
            dashboard=self.dashboard,
            user=self.user,
            metadata=self.metadata,
            plan={
                'actions': actions,
                'message': '',
            },
        )

    def test_update_pie_to_bar(self):
        widget = self.dashboard.widgets.get(
            settings__source='missing_values'
        )

        count = self._apply([{
            'action': 'update',
            'widget_id': widget.pk,
            'chart_type': 'bar',
        }])

        widget.refresh_from_db()
        self.assertEqual(count, 1)
        self.assertEqual(widget.widget_type, 'bar')

    def test_delete_widget(self):
        widget = self.dashboard.widgets.first()

        count = self._apply([{
            'action': 'delete',
            'widget_id': widget.pk,
        }])

        self.assertEqual(count, 1)
        self.assertFalse(
            DashboardWidget.objects.filter(
                pk=widget.pk
            ).exists()
        )

    def test_add_widget_with_valid_columns(self):
        count = self._apply([{
            'action': 'add',
            'title': 'المبيعات الشهرية',
            'widget_type': 'chart',
            'chart_type': 'line',
            'sheet_name': 'المبيعات',
            'x_column': 'الشهر',
            'y_column': 'القيمة',
            'aggregation': 'sum',
        }])

        widget = self.dashboard.widgets.get(
            title='المبيعات الشهرية'
        )
        self.assertEqual(count, 1)
        self.assertEqual(widget.widget_type, 'line')
        self.assertEqual(
            widget.settings['sheet_name'],
            'المبيعات',
        )

    def test_rejects_unknown_column(self):
        with self.assertRaises(
            DashboardPlanValidationError
        ):
            self._apply([{
                'action': 'add',
                'title': 'مخطط خاطئ',
                'widget_type': 'chart',
                'chart_type': 'bar',
                'sheet_name': 'المبيعات',
                'x_column': 'عمود غير موجود',
                'y_column': 'القيمة',
                'aggregation': 'sum',
            }])

    def test_rejects_widget_from_another_dashboard(self):
        other_user = get_user_model().objects.create_user(
            username='other-ai-user',
            password='strong-test-password',
        )
        other_dashboard = Dashboard.objects.create(
            owner=other_user,
            name='لوحة أخرى',
        )
        foreign_widget = DashboardWidget.objects.create(
            dashboard=other_dashboard,
            title='عنصر خاص',
            widget_type='pie',
        )

        with self.assertRaises(
            DashboardPlanValidationError
        ):
            self._apply([{
                'action': 'delete',
                'widget_id': foreign_widget.pk,
            }])

        self.assertTrue(
            DashboardWidget.objects.filter(
                pk=foreign_widget.pk
            ).exists()
        )

    def test_rejects_invalid_chart_type(self):
        widget = self.dashboard.widgets.first()

        with self.assertRaises(
            DashboardPlanValidationError
        ):
            self._apply([{
                'action': 'update',
                'widget_id': widget.pk,
                'chart_type': 'scatter',
            }])

    def test_invalid_plan_json_shape_is_rejected(self):
        with self.assertRaises(
            DashboardPlanValidationError
        ):
            apply_dashboard_plan(
                dashboard=self.dashboard,
                user=self.user,
                metadata=self.metadata,
                plan={'actions': 'not-a-list'},
            )

    def test_transaction_rolls_back_all_actions(self):
        widget = self.dashboard.widgets.first()
        original_title = widget.title

        with self.assertRaises(
            DashboardPlanValidationError
        ):
            self._apply([
                {
                    'action': 'update',
                    'widget_id': widget.pk,
                    'title': 'عنوان مؤقت',
                },
                {
                    'action': 'add',
                    'title': 'مخطط خاطئ',
                    'widget_type': 'chart',
                    'chart_type': 'bar',
                    'sheet_name': 'المبيعات',
                    'x_column': 'غير موجود',
                    'y_column': 'القيمة',
                    'aggregation': 'sum',
                },
            ])

        widget.refresh_from_db()
        self.assertEqual(widget.title, original_title)

    def test_missing_api_key_returns_safe_arabic_error(self):
        with patch.dict(
            os.environ,
            {'OPENAI_API_KEY': ''},
        ):
            with self.assertRaisesRegex(
                DashboardAIUnavailableError,
                'غير مهيأة',
            ):
                request_dashboard_plan(
                    user_request='غيّر المخطط',
                    metadata=self.metadata,
                )

    def test_update_single_widget_colors(self):
        widget = self.dashboard.widgets.first()

        count = self._apply([{
            'action': 'update',
            'widget_id': widget.pk,
            'colors': ['#ec4899', '#facc15'],
        }])

        widget.refresh_from_db()
        self.assertEqual(count, 1)
        self.assertEqual(
            widget.settings['colors'],
            ['#ec4899', '#facc15'],
        )

    def test_update_all_dashboard_widget_colors(self):
        widgets = list(self.dashboard.widgets.all())
        actions = [
            {
                'action': 'update',
                'widget_id': widget.pk,
                'colors': ['#ec4899', '#facc15'],
            }
            for widget in widgets
        ]

        count = self._apply(actions)

        self.assertEqual(count, len(widgets))
        for widget in DashboardWidget.objects.filter(
            dashboard=self.dashboard
        ):
            self.assertEqual(
                widget.settings['colors'],
                ['#ec4899', '#facc15'],
            )

    def test_rejects_invalid_hex_color(self):
        widget = self.dashboard.widgets.first()

        with self.assertRaises(
            DashboardPlanValidationError
        ):
            self._apply([{
                'action': 'update',
                'widget_id': widget.pk,
                'colors': ['pink', 'javascript:alert(1)'],
            }])

        widget.refresh_from_db()
        self.assertNotIn(
            'colors',
            widget.settings,
        )

    def test_rejects_more_than_ten_colors(self):
        widget = self.dashboard.widgets.first()

        with self.assertRaises(
            DashboardPlanValidationError
        ):
            self._apply([{
                'action': 'update',
                'widget_id': widget.pk,
                'colors': ['#ec4899'] * 11,
            }])

    def test_widget_colors_survive_reload(self):
        widget = self.dashboard.widgets.first()
        self._apply([{
            'action': 'update',
            'widget_id': widget.pk,
            'colors': ['#3b82f6', '#22c55e'],
        }])

        reloaded_dashboard = get_or_create_dataset_dashboard(
            self.dataset,
            has_numeric_columns=True,
        )
        reloaded_widget = reloaded_dashboard.widgets.get(
            pk=widget.pk
        )

        self.assertEqual(
            reloaded_widget.settings['colors'],
            ['#3b82f6', '#22c55e'],
        )

    def test_cannot_change_colors_of_another_users_widget(self):
        other_user = get_user_model().objects.create_user(
            username='color-owner',
            password='strong-test-password',
        )
        other_dashboard = Dashboard.objects.create(
            owner=other_user,
            name='لوحة ألوان خاصة',
        )
        foreign_widget = DashboardWidget.objects.create(
            dashboard=other_dashboard,
            title='مخطط خاص',
            widget_type='pie',
        )

        with self.assertRaises(
            DashboardPlanValidationError
        ):
            self._apply([{
                'action': 'update',
                'widget_id': foreign_widget.pk,
                'colors': ['#ec4899'],
            }])

        foreign_widget.refresh_from_db()
        self.assertNotIn(
            'colors',
            foreign_widget.settings,
        )
