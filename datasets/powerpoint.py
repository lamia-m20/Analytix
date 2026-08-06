import io
import re
from collections import Counter

import matplotlib.pyplot as plt
from arabic_reshaper import reshape
from bidi.algorithm import get_display
from matplotlib import font_manager
from matplotlib.font_manager import FontProperties
from matplotlib.ticker import FuncFormatter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .exporters import TOP_RESULTS, get_arabic_font_path

MAX_POWERPOINT_SLIDES = 40
MAX_POWERPOINT_CHARTS = 12
PPTX_CONTENT_TYPE = (
    'application/vnd.openxmlformats-officedocument.presentationml.presentation'
)

NAVY = RGBColor(23, 50, 74)
DEEP_BLUE = RGBColor(32, 68, 93)
BLUE = RGBColor(72, 180, 242)
SKY = RGBColor(57, 198, 216)
TEAL = RGBColor(110, 215, 200)
GREEN = RGBColor(65, 211, 154)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
WHITE = RGBColor(255, 255, 255)
SLATE = RGBColor(199, 217, 230)
PALE = RGBColor(42, 83, 108)
BACKGROUND = RGBColor(23, 50, 74)
PANEL = RGBColor(32, 68, 93)
PANEL_LIGHT = RGBColor(42, 83, 108)
PANEL_ALT = RGBColor(42, 83, 108)
PANEL_BORDER = RGBColor(59, 106, 130)
FONT_NAME = 'Tahoma'
PPT_CHART_PALETTE = ('#39C6D8', '#6ED7C8', '#48B4F2', '#41D39A', '#F59E0B')
ARABIC_CHART_PATTERN = re.compile(r'[\u0600-\u06FF]')

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
CONTENT_LEFT = Inches(.62)
CONTENT_RIGHT = Inches(.62)
HEADER_HEIGHT = Inches(.18)
FOOTER_HEIGHT = Inches(.38)
TITLE_SIZE = 29
SUBTITLE_SIZE = 17
BODY_SIZE = 17
KPI_SIZE = 30
SMALL_SIZE = 14


def prepare_arabic_for_chart(text):
    """Shape bidi text only for rasterized Matplotlib chart content."""
    text = str(text)
    if not ARABIC_CHART_PATTERN.search(text):
        return text
    protected = []

    def protect(match):
        protected.append(match.group(0))
        return chr(0xE000 + len(protected) - 1)

    safe_text = re.sub(r'[A-Za-z][A-Za-z0-9_.-]*|\d+(?:[.,]\d+)*%?', protect, text)
    display = get_display(reshape(safe_text))
    for index, token in enumerate(protected):
        display = display.replace(chr(0xE000 + index), token)
    return display


def _chart_font():
    """Use the project/PDF Arabic font without a platform-specific path."""
    try:
        path = get_arabic_font_path()
        font_manager.fontManager.addfont(str(path))
        return FontProperties(fname=str(path))
    except Exception:
        return FontProperties(family='DejaVu Sans')


def _set_text(frame, text, size=20, bold=False, color=WHITE, align=PP_ALIGN.RIGHT):
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(.08); frame.margin_right = Inches(.08)
    frame.margin_top = Inches(.04); frame.margin_bottom = Inches(.04)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    if align == PP_ALIGN.RIGHT:
        paragraph._p.get_or_add_pPr().set('rtl', '1')
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _textbox(slide, text, left, top, width, height, size=18, bold=False,
             color=WHITE, align=PP_ALIGN.RIGHT):
    box = slide.shapes.add_textbox(left, top, width, height)
    _set_text(box.text_frame, text, size, bold, color, align)
    return box


def apply_arabic_text_style(frame, text, size=BODY_SIZE, bold=False, color=WHITE):
    _set_text(frame, text, size, bold, color, PP_ALIGN.RIGHT)
    return frame


def add_modern_header(slide, prs):
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.08))
    banner.name = 'ANALYTIX_HEADER'
    banner.fill.solid(); banner.fill.fore_color.rgb = NAVY; banner.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.9), Inches(.08), Inches(1.7), Inches(.025))
    accent.fill.solid(); accent.fill.fore_color.rgb = SKY; accent.line.fill.background()


def add_title(slide, title):
    box = _textbox(slide, title, Inches(.72), Inches(.3), Inches(11.85), Inches(.62), TITLE_SIZE, True, WHITE)
    box.name = 'SLIDE_TITLE'
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.55), Inches(.94), Inches(1), Inches(.035))
    rule.fill.solid(); rule.fill.fore_color.rgb = SKY; rule.line.fill.background()
    return box


def add_subtitle(slide, text, top=1.05):
    box = _textbox(slide, text, Inches(.75), Inches(top), Inches(11.8), Inches(.42), SUBTITLE_SIZE, False, SLATE)
    box.name = 'SLIDE_SUBTITLE'
    return box


def add_page_number(slide, number):
    box = _textbox(slide, str(number), Inches(12.1), Inches(7.08), Inches(.48), Inches(.22), 10, False, SLATE, PP_ALIGN.CENTER)
    box.name = 'PAGE_NUMBER'


def add_modern_footer(slide, number):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.62), Inches(6.96), Inches(12.05), Inches(.015))
    line.fill.solid(); line.fill.fore_color.rgb = PANEL_BORDER; line.line.fill.background()
    footer = _textbox(slide, 'Analytix • Data Intelligence', Inches(.62), Inches(7.05), Inches(2.6), Inches(.22), 10, True, SLATE, PP_ALIGN.LEFT)
    footer.name = 'ANALYTIX_FOOTER'
    add_page_number(slide, number)


def add_badge(slide, text, left, top, color=BLUE):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.5), Inches(.38))
    badge.name = 'BADGE'; badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
    _set_text(badge.text_frame, text, 11, True, WHITE, PP_ALIGN.CENTER)
    return badge


def add_section_card(slide, left, top, width, height, title=None, border=SKY):
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(.04), top + Inches(.06), width, height)
    shadow.fill.solid(); shadow.fill.fore_color.rgb = RGBColor(2, 6, 15); shadow.line.fill.background()
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.name = 'SECTION_CARD'; panel.fill.solid(); panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = PANEL_BORDER; panel.line.width = Pt(.6)
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + width - Inches(.045), top, Inches(.045), height)
    marker.fill.solid(); marker.fill.fore_color.rgb = border; marker.line.fill.background()
    if title:
        _textbox(slide, title, left + Inches(.2), top + Inches(.08), width - Inches(.38), Inches(.38), 15, True, border)
    return panel


def add_bullet_panel(slide, title, items, left, top, width, height, max_items=5, color=BLUE):
    add_section_card(slide, left, top, width, height, title, color)
    box = slide.shapes.add_textbox(left + Inches(.24), top + Inches(.48), width - Inches(.48), height - Inches(.62))
    box.name = 'BULLET_PANEL'
    frame = box.text_frame; frame.clear(); frame.word_wrap = True; frame.margin_right = 0
    for index, item in enumerate(items[:max_items]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f'• {str(item)[:150]}'
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph._p.get_or_add_pPr().set('rtl', '1')
        paragraph.font.name = FONT_NAME; paragraph.font.size = Pt(BODY_SIZE); paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(7)
    return box


def add_insight_panel(slide, text):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.72), Inches(6.05), Inches(11.95), Inches(.72))
    panel.name = 'INSIGHT_PANEL'; panel.fill.solid(); panel.fill.fore_color.rgb = PANEL_ALT
    panel.line.color.rgb = PANEL_BORDER; panel.line.width = Pt(.5)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.56), Inches(6.16), Inches(.035), Inches(.48))
    line.fill.solid(); line.fill.fore_color.rgb = SKY; line.line.fill.background()
    _textbox(slide, 'القراءة الرئيسية', Inches(10.25), Inches(6.18), Inches(2.05), Inches(.3), 14, True, SKY)
    _textbox(slide, text[:250], Inches(.98), Inches(6.18), Inches(9.05), Inches(.34), 13, False, WHITE)
    return panel


def add_chart_slide_modern(slide, chart, insight='', metrics=None):
    metrics = metrics or []
    for index, (label, value, color) in enumerate(metrics[:4]):
        left = Inches(.72 + index * 3.04)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.12), Inches(2.82), Inches(.76))
        card.name = 'CHART_KPI'; card.fill.solid(); card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = PANEL_BORDER; card.line.width = Pt(.5)
        _textbox(slide, str(value), left + Inches(.12), Inches(1.2), Inches(1.05), Inches(.4), 19, True, color, PP_ALIGN.LEFT)
        _textbox(slide, label, left + Inches(1.1), Inches(1.24), Inches(1.52), Inches(.3), 11, False, SLATE)
    chart_width = Inches(11.95)
    chart_left = (SLIDE_WIDTH - chart_width) // 2
    picture = slide.shapes.add_picture(
        chart, chart_left, Inches(1.98),
        width=chart_width, height=Inches(3.92),
    )
    picture.name = 'CHART_PICTURE'
    add_insight_panel(slide, insight or 'يعرض الرسم أبرز القيم المحسوبة من البيانات، ويُوصى بالبدء بالفئة الأعلى أثرًا عند المراجعة.')
    return picture


def add_table_panel(slide, title, rows, left, top, width, height):
    """Compact text-table panel for the few top rows used in presentations."""
    lines = [' | '.join(str(value) for value in row) for row in rows[:8]]
    return add_bullet_panel(slide, title, lines, left, top, width, height, max_items=8)


def add_clean_table(slide, rows, left, top, width, row_height=Inches(.72)):
    severity_colors = {'مرتفع': RED, 'متوسط': AMBER, 'منخفض': TEAL}
    for index, row in enumerate(rows[:5]):
        y = top + index * row_height
        if index:
            rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, width, Inches(.012))
            rule.fill.solid(); rule.fill.fore_color.rgb = PANEL_BORDER; rule.line.fill.background()
        _textbox(slide, row[0], left + int(width * .38), y + Inches(.12), int(width * .58), Inches(.42), 16, index == 0, WHITE)
        _textbox(slide, row[1], left + int(width * .18), y + Inches(.12), int(width * .18), Inches(.42), 15, True, BLUE, PP_ALIGN.CENTER)
        _textbox(slide, row[2], left, y + Inches(.12), int(width * .16), Inches(.42), 13, False, SLATE, PP_ALIGN.CENTER)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left + Inches(.05), y + Inches(.24), Inches(.11), Inches(.11)
        )
        dot.name = 'SEVERITY_DOT'
        dot.fill.solid(); dot.fill.fore_color.rgb = severity_colors.get(row[2], SLATE)
        dot.line.fill.background()


def add_priority_columns(slide, columns):
    colors = [RED, AMBER, BLUE]
    for index, (title, items) in enumerate(columns):
        left = Inches(.72 + index * 4.14)
        _textbox(slide, title, left, Inches(1.62), Inches(3.72), Inches(.48), 19, True, colors[index])
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.18), Inches(3.72), Inches(.035))
        rule.fill.solid(); rule.fill.fore_color.rgb = colors[index]; rule.line.fill.background()
        for item_index, item in enumerate(items[:3]):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(3.48), Inches(2.55 + item_index * 1.22), Inches(.12), Inches(.12))
            dot.fill.solid(); dot.fill.fore_color.rgb = colors[index]; dot.line.fill.background()
            _textbox(slide, str(item)[:150], left, Inches(2.43 + item_index * 1.22), Inches(3.3), Inches(.82), 16, False, WHITE)


def _base_slide(prs, title, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid(); background.fore_color.rgb = BACKGROUND
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.1), Inches(.35), Inches(7.1), Inches(6.7))
    glow.name = 'SOFT_GLOW'; glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(31, 78, 102)
    glow.fill.transparency = 58; glow.line.fill.background()
    add_modern_header(slide, prs); add_title(slide, title); add_modern_footer(slide, number)
    return slide


def _bullet_list(slide, items, top=1.2, max_items=8, color=WHITE):
    box = slide.shapes.add_textbox(Inches(.75), Inches(top), Inches(11.8), Inches(5.45))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    for index, item in enumerate(items[:max_items]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f'• {item}'
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.font.name = FONT_NAME; paragraph.font.size = Pt(18); paragraph.font.color.rgb = color
        paragraph.space_after = Pt(10)
    return box


def add_kpi_card_modern(slide, label, value, left, top, color=BLUE, width=Inches(3.72)):
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(.035), top + Inches(.055), width, Inches(1.24))
    shadow.fill.solid(); shadow.fill.fore_color.rgb = RGBColor(2, 6, 15); shadow.line.fill.background()
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.24))
    shape.name = 'KPI_CARD'; shape.fill.solid(); shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = PANEL_BORDER; shape.line.width = Pt(.55)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(.045))
    accent.fill.solid(); accent.fill.fore_color.rgb = color; accent.line.fill.background()
    text_width = width - Inches(.54)
    _textbox(slide, value, left + Inches(.16), top + Inches(.12), text_width, Inches(.58), KPI_SIZE, True, color, PP_ALIGN.CENTER)
    _textbox(slide, label, left + Inches(.16), top + Inches(.73), text_width, Inches(.32), SMALL_SIZE, False, SLATE)


_kpi_card = add_kpi_card_modern


def _dashboard_chart(labels, values, kind='bar', colors_list=None):
    """Render an Arabic-safe chart used only by PowerPoint."""
    buffer = io.BytesIO()
    figure = None
    try:
        palette = colors_list or PPT_CHART_PALETTE
        chart_font = _chart_font()
        display_labels = [prepare_arabic_for_chart(value) for value in labels]
        figure, axis = plt.subplots(figsize=(12, 4.4), facecolor='#20445D')
        axis.set_facecolor('#20445D')
        figure.patch.set_alpha(1)
        axis.tick_params(colors='#94A3B8', labelsize=10, length=0)
        axis.set_axisbelow(True)
        for spine in axis.spines.values():
            spine.set_visible(False)
        if kind in ('pie', 'donut'):
            wedges, _, autotexts = axis.pie(
                values, colors=[palette[i % len(palette)] for i in range(len(values))],
                startangle=90, autopct='%1.0f%%', pctdistance=.78,
                wedgeprops={'width': .38, 'edgecolor': '#20445D', 'linewidth': 3},
            )
            for text in autotexts:
                text.set_color('#F8FAFC'); text.set_fontweight('bold'); text.set_fontsize(10)
                text.set_fontproperties(chart_font)
            axis.legend(wedges, display_labels, frameon=False, labelcolor='#CBD5E1',
                        loc='center left', bbox_to_anchor=(-.05, .5), fontsize=9, prop=chart_font)
            axis.axis('equal')
        else:
            positions = range(len(values))
            bars = axis.bar(positions, values, color=palette, width=.55)
            axis.set_xticks(list(positions))
            axis.set_xticklabels(
                display_labels,
                rotation=24 if len(values) > 5 else 0,
                ha='right' if len(values) > 5 else 'center',
            )
            axis.grid(axis='y', color='#3B6A82', linewidth=.6, alpha=.5)
            labels_drawn = axis.bar_label(bars, fmt='{:,.0f}', padding=5, color='#E2E8F0', fontsize=10, fontweight='bold')
            for label in labels_drawn: label.set_fontproperties(chart_font)
        if kind not in ('pie', 'donut'):
            axis.yaxis.set_major_formatter(FuncFormatter(lambda value, _: f'{value:,.0f}'))
            for label in (*axis.get_xticklabels(), *axis.get_yticklabels()):
                label.set_fontproperties(chart_font)
        figure.subplots_adjust(
            left=.07, right=.97,
            bottom=.25 if len(values) > 5 and kind not in ('pie', 'donut') else .13,
            top=.95,
        )
        figure.savefig(buffer, format='png', dpi=300, facecolor='#20445D')
        buffer.seek(0)
        return buffer
    except Exception:
        buffer.close()
        return None
    finally:
        if figure is not None:
            plt.close(figure)


def _dashboard_widget_chart(widget, report):
    try:
        sheet = next(item for item in report['sheets'] if item['name'] == widget['sheet'])
        frame = sheet['_frame']
        grouped = frame.groupby(widget['x_column'], dropna=False)[widget['y_column']]
        aggregation = widget.get('aggregation')
        values = grouped.sum() if aggregation == 'sum' else grouped.mean() if aggregation in ('average', 'mean') else grouped.count()
        values = values.head(TOP_RESULTS)
        return _dashboard_chart(
            [str(value) for value in values.index], list(values.values),
            'bar',
        )
    except Exception:
        return None


def _chart_slide(prs, title, labels, values, number, kind='bar', palette=None):
    short_labels = [str(value)[:24] + ('…' if len(str(value)) > 24 else '') for value in labels[:8]]
    chart_colors = list(palette or PPT_CHART_PALETTE)
    if kind == 'bar' and values:
        maximum_index = max(range(len(values[:8])), key=lambda index: values[index])
        chart_colors = ['#2563EB'] * len(values[:8])
        chart_colors[maximum_index] = '#06B6D4'
    chart = _dashboard_chart(short_labels, values[:8], 'bar', chart_colors)
    if not chart:
        return False
    try:
        slide = _base_slide(prs, title, number)
        if values:
            maximum_index = max(range(len(values[:8])), key=lambda index: values[index])
            insight = f'أعلى قيمة ظهرت في {labels[maximum_index]} وبلغت {values[maximum_index]:,.2f}.'
        else:
            insight = 'لا توجد قيم كافية لاستخلاص مقارنة رقمية.'
        reading = f'{insight} يوصى بمراجعة الفئة الأعلى أولًا ثم مقارنتها ببقية الفئات.'
        metric_values = values[:8]
        metrics = [
            ('عدد الفئات', len(short_labels), BLUE),
            ('أعلى قيمة', f'{max(metric_values):,.1f}', SKY),
            ('المتوسط', f'{sum(metric_values) / len(metric_values):,.1f}', TEAL),
            ('الإجمالي', f'{sum(metric_values):,.1f}', AMBER),
        ] if metric_values else []
        add_chart_slide_modern(slide, chart, reading, metrics)
        return True
    finally:
        chart.close()


def _priority(issue):
    return {'مرتفع': 'أولوية مرتفعة', 'متوسط': 'أولوية متوسطة'}.get(issue, 'أولوية منخفضة')


def add_dashboard_sidebar(slide, title, active='النظرة العامة'):
    sidebar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.05), Inches(.22), Inches(2.05), Inches(6.86))
    sidebar.name = 'DASHBOARD_SIDEBAR'; sidebar.fill.solid(); sidebar.fill.fore_color.rgb = PANEL
    sidebar.line.color.rgb = PANEL_BORDER; sidebar.line.width = Pt(.6)
    _textbox(slide, 'Analytix', Inches(11.32), Inches(.52), Inches(1.5), Inches(.38), 19, True, SKY)
    _textbox(slide, title[:48], Inches(11.28), Inches(1.02), Inches(1.52), Inches(.75), 13, True, WHITE)
    sections = ('النظرة العامة', 'جودة البيانات', 'المشكلات', 'الرسوم', 'التوصيات')
    for index, section in enumerate(sections):
        top = Inches(2.0 + index * .72)
        if section == active:
            marker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.22), top - Inches(.04), Inches(1.7), Inches(.48))
            marker.fill.solid(); marker.fill.fore_color.rgb = PANEL_LIGHT
            marker.line.fill.background()
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.58), top + Inches(.09), Inches(.1), Inches(.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = SKY if section == active else PANEL_BORDER; dot.line.fill.background()
        _textbox(slide, section, Inches(11.37), top, Inches(1.08), Inches(.28), 11, section == active, WHITE if section == active else SLATE)


def add_dashboard_kpi(slide, label, value, left, top, color=SKY, icon='•', width=Inches(1.58)):
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(.025), top + Inches(.04), width, Inches(.88))
    shadow.fill.solid(); shadow.fill.fore_color.rgb = RGBColor(17, 43, 62); shadow.line.fill.background()
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(.88))
    card.name = 'KPI_CARD'; card.fill.solid(); card.fill.fore_color.rgb = PANEL_LIGHT
    card.line.color.rgb = PANEL_BORDER; card.line.width = Pt(.45)
    _textbox(slide, icon, left + Inches(.1), top + Inches(.14), Inches(.3), Inches(.28), 13, True, color, PP_ALIGN.CENTER)
    _textbox(slide, str(value), left + Inches(.43), top + Inches(.08), width - Inches(.54), Inches(.38), 23, True, WHITE)
    _textbox(slide, label, left + Inches(.18), top + Inches(.53), width - Inches(.34), Inches(.22), 10, False, SLATE)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(.14), top + Inches(.8), width - Inches(.28), Inches(.025))
    accent.fill.solid(); accent.fill.fore_color.rgb = color; accent.line.fill.background()
    return card


def _dashboard_gauge(value):
    buffer = io.BytesIO(); figure = None
    try:
        value = max(0, min(100, float(value)))
        chart_font = _chart_font()
        figure, axis = plt.subplots(figsize=(3.2, 2.5), facecolor='#20445D')
        axis.set_facecolor('#20445D')
        axis.pie([value, 100 - value], startangle=90, counterclock=False,
                 colors=['#41D39A', '#3B6A82'],
                 wedgeprops={'width': .24, 'edgecolor': '#20445D', 'linewidth': 2})
        axis.text(0, .02, prepare_arabic_for_chart(f'{value:.0f}%'), ha='center', va='center', color='#F4FAFF', fontsize=23, fontweight='bold', fontproperties=chart_font)
        axis.text(0, -.3, prepare_arabic_for_chart('درجة الجودة'), ha='center', color='#C7D9E6', fontsize=9, fontproperties=chart_font)
        axis.axis('equal'); axis.axis('off'); figure.tight_layout(pad=.1)
        figure.savefig(buffer, format='png', dpi=300, facecolor='#20445D', bbox_inches='tight')
        buffer.seek(0); return buffer
    finally:
        if figure is not None: plt.close(figure)


def add_image_panel(slide, image, left, top, width, height, name):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid(); panel.fill.fore_color.rgb = PANEL; panel.line.color.rgb = PANEL_BORDER; panel.line.width = Pt(.45)
    picture = slide.shapes.add_picture(image, left + Inches(.08), top + Inches(.08), width=width - Inches(.16), height=height - Inches(.16))
    picture.name = name
    return picture


def build_powerpoint_report(dataset, report):
    """Build a bounded, visual presentation from the shared report_data."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    chart_count = 0

    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = NAVY
    cyan_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(.18), SLIDE_HEIGHT)
    cyan_block.fill.solid(); cyan_block.fill.fore_color.rgb = SKY; cyan_block.line.fill.background()
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.85), Inches(1.34), Inches(2.25), Inches(.045))
    line.fill.solid(); line.fill.fore_color.rgb = SKY; line.line.fill.background()
    _textbox(slide, 'Analytix', Inches(8.75), Inches(.55), Inches(3.35), Inches(.45), 16, True, SKY)
    _textbox(slide, 'تقرير تحليل البيانات الشامل', Inches(3.1), Inches(1.55), Inches(9), Inches(1.15), 34, True, WHITE)
    details = (
        f'{report["dataset"]["filename"]}\n'
        f'المستخدم: {report["user"]["username"]}\n'
        f'Dataset ID: {report["dataset"]["id"]}\n'
        f'{report["generated_at"].strftime("%Y-%m-%d %H:%M")}'
    )
    _textbox(slide, details, Inches(5.25), Inches(4.55), Inches(6.85), Inches(1.75), 17, False, WHITE)

    summary = report['summary']
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BACKGROUND
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.2), Inches(.3), Inches(8), Inches(6.7))
    glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(45, 100, 123); glow.fill.transparency = 62; glow.line.fill.background()
    add_dashboard_sidebar(slide, 'تقرير تحليل البيانات', 'النظرة العامة')
    add_modern_header(slide, prs); add_modern_footer(slide, len(prs.slides))
    _textbox(slide, 'الملخص التنفيذي | لوحة المؤشرات', Inches(.55), Inches(.28), Inches(10.15), Inches(.5), 28, True, WHITE)
    _textbox(slide, report['dataset']['filename'], Inches(.55), Inches(.78), Inches(10.15), Inches(.28), 12, False, SLATE)
    cards = [
        ('عدد الأوراق', summary['sheet_count'], SKY, '▦'), ('عدد الصفوف', f'{summary["rows"]:,}', BLUE, '≡'),
        ('عدد الأعمدة', summary['columns'], TEAL, '▥'), ('درجة الجودة', f'{summary["quality_score"]}%', GREEN, '✓'),
        ('القيم الفارغة', f'{summary["missing_cells"]:,}', AMBER, '!'), ('المكررة', f'{summary["duplicates"]:,}', RED, '↻'),
    ]
    for index, (label, value, color, icon) in enumerate(cards):
        add_dashboard_kpi(slide, label, value, Inches(.55 + index * 1.72), Inches(1.12), color, icon)
    sheet_labels = [item['name'] for item in report['sheets'][:8]]
    quality_values = [item['quality_score'] for item in report['sheets'][:8]]
    bar = _dashboard_chart(sheet_labels, quality_values, 'bar', ['#39C6D8'] * max(1, len(quality_values)))
    donut = _dashboard_chart(['مكتملة', 'فارغة'], [summary['completion_rate'], max(0, 100 - summary['completion_rate'])], 'donut', ['#6ED7C8', '#3B6A82'])
    gauge = _dashboard_gauge(summary['quality_score'])
    try:
        if bar: add_image_panel(slide, bar, Inches(.55), Inches(2.18), Inches(4.65), Inches(2.72), 'EXECUTIVE_BAR')
        if donut: add_image_panel(slide, donut, Inches(5.38), Inches(2.18), Inches(2.55), Inches(2.72), 'EXECUTIVE_DONUT')
        if gauge: add_image_panel(slide, gauge, Inches(8.1), Inches(2.18), Inches(2.55), Inches(2.72), 'EXECUTIVE_GAUGE')
    finally:
        if bar: bar.close()
        if donut: donut.close()
        if gauge: gauge.close()
    results_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.55), Inches(5.1), Inches(10.1), Inches(1.58))
    results_panel.fill.solid(); results_panel.fill.fore_color.rgb = PANEL; results_panel.line.color.rgb = PANEL_BORDER; results_panel.line.width = Pt(.45)
    _textbox(slide, 'أبرز النتائج', Inches(8.35), Inches(5.3), Inches(1.95), Inches(.3), 16, True, SKY)
    insights = report['executive_summary']['insights'][:4]
    for index, insight in enumerate(insights):
        column = index % 2; row = index // 2
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.9 - column * 4.9), Inches(5.82 + row * .42), Inches(.09), Inches(.09))
        dot.fill.solid(); dot.fill.fore_color.rgb = TEAL; dot.line.fill.background()
        _textbox(slide, str(insight)[:92], Inches(5.35 - column * 4.8), Inches(5.7 + row * .42), Inches(4.38), Inches(.34), 11, False, WHITE)

    slide = _base_slide(prs, 'مؤشرات الجودة', len(prs.slides) + 1)
    high_issues = sum(1 for sheet in report['sheets'] for issue in sheet['issues'] if issue['severity'] == 'مرتفع')
    duplicate_rate = round(summary['duplicates'] / summary['rows'] * 100, 2) if summary['rows'] else 0
    quality_cards = [
        ('اكتمال البيانات', f'{summary["completion_rate"]}%'),
        ('نسبة التكرار', f'{duplicate_rate}%'), ('مشكلات مرتفعة', high_issues),
        ('عدد الأوراق', summary['sheet_count']), ('الجودة العامة', f'{summary["quality_score"]}/100'),
    ]
    for index, (label, value) in enumerate(quality_cards):
        _kpi_card(slide, label, value, Inches(.72 + (index % 3) * 4.16), Inches(1.4 + (index // 3) * 1.62), [BLUE, RED, TEAL][index % 3])
    add_bullet_panel(slide, 'قراءة سريعة', [f'درجة الجودة العامة {summary["quality_score"]} من 100.', f'بلغ اكتمال البيانات {summary["completion_rate"]}%.'], Inches(.72), Inches(4.65), Inches(12), Inches(1.75), max_items=2, color=BLUE)

    issues = report['executive_summary']['issues']
    slide = _base_slide(prs, 'أهم المشكلات', len(prs.slides) + 1)
    if issues:
        rows = []
        for issue in issues[:5]:
            number = issue['issue'].split(':', 1)[-1].strip()[:30]
            rows.append((f'{issue["sheet"]}: {issue["issue"].split(":", 1)[0]}', number, issue['severity']))
        add_clean_table(slide, rows, Inches(.72), Inches(1.42), Inches(12), row_height=Inches(.94))
    else:
        add_bullet_panel(slide, 'حالة الجودة', ['لم تُرصد مشكلات بارزة وفق قواعد التقرير.'], Inches(.72), Inches(1.45), Inches(12), Inches(2), color=TEAL)

    recommendations = report['executive_summary']['recommendations']
    recommendation_lines = []
    for index, recommendation in enumerate(recommendations):
        severity = issues[index]['severity'] if index < len(issues) else 'منخفض'
        recommendation_lines.append(f'{_priority(severity)}: {recommendation}')
    recommendation_lines = recommendation_lines or ['أولوية منخفضة: استمر في المراجعة الدورية لجودة البيانات.']
    priority_groups = {'مرتفعة': [], 'متوسطة': [], 'منخفضة': []}
    for recommendation in recommendation_lines:
        target = 'مرتفعة' if 'مرتفعة' in recommendation else 'متوسطة' if 'متوسطة' in recommendation else 'منخفضة'
        priority_groups[target].append(recommendation.partition(':')[2].strip())
    slide = _base_slide(prs, 'خارطة الأولويات', len(prs.slides) + 1)
    add_subtitle(slide, 'التوصيات موزعة حسب أولوية التنفيذ')
    add_priority_columns(slide, [(name, priority_groups[name]) for name in ('مرتفعة', 'متوسطة', 'منخفضة')])

    chart_specs = [
        ('جودة الأوراق', [x['name'] for x in report['sheets']], [x['quality_score'] for x in report['sheets']], 'bar'),
        ('القيم الفارغة', [x['name'] for x in report['sheets']], [x['missing_cells'] for x in report['sheets']], 'bar'),
        ('الصفوف المكررة', [x['name'] for x in report['sheets']], [x['duplicates'] for x in report['sheets']], 'bar'),
    ]
    type_counts = Counter(
        'رقمي' if 'int' in c['type'] or 'float' in c['type'] else
        'تاريخي' if 'date' in c['type'] else 'منطقي' if 'bool' in c['type'] else 'نصي'
        for sheet in report['sheets'] for c in sheet['columns_info']
    )
    chart_specs.insert(2, ('أنواع الأعمدة', list(type_counts), list(type_counts.values()), 'pie'))
    outliers = sorted(
        [(f'{sheet["name"]}: {item["column"]}', item['count']) for sheet in report['sheets'] for item in sheet['outliers'] if item['count']],
        key=lambda value: value[1], reverse=True,
    )[:TOP_RESULTS]
    if outliers: chart_specs.append(('القيم الشاذة', [x[0] for x in outliers], [x[1] for x in outliers], 'bar'))
    correlations = sorted(
        [(f'{sheet["name"]}: {item["column_1"]} / {item["column_2"]}', abs(item['value'])) for sheet in report['sheets'] for item in sheet['correlations']],
        key=lambda value: value[1], reverse=True,
    )[:TOP_RESULTS]
    if correlations: chart_specs.append(('أقوى الارتباطات', [x[0] for x in correlations], [x[1] for x in correlations], 'bar'))
    for title, labels, values, kind in chart_specs:
        if len(prs.slides) >= MAX_POWERPOINT_SLIDES - 1 or chart_count >= MAX_POWERPOINT_CHARTS: break
        if labels and _chart_slide(prs, title, labels, values, len(prs.slides) + 1, kind): chart_count += 1

    for widget in report['widgets']:
        if len(prs.slides) >= MAX_POWERPOINT_SLIDES - 1: break
        if widget['type'] in ('bar', 'line', 'pie') and chart_count < MAX_POWERPOINT_CHARTS:
            image = _dashboard_widget_chart(widget, report)
            if image:
                try:
                    slide = _base_slide(prs, widget['title'], len(prs.slides) + 1)
                    add_chart_slide_modern(slide, image, 'يعرض هذا المخطط العنصر المحفوظ في لوحة المعلومات الحالية وأبرز المقارنة الناتجة عنه.')
                    chart_count += 1
                finally: image.close()
        elif widget['type'] in ('metric', 'table'):
            slide = _base_slide(prs, widget['title'], len(prs.slides) + 1)
            add_bullet_panel(slide, 'تفاصيل العنصر', [f'الورقة: {widget["sheet"] or "غير محددة"}', f'النوع: {widget["type"]}', f'التجميع: {widget["aggregation"] or "بدون"}'], Inches(.72), Inches(1.35), Inches(12), Inches(3.3))

    for sheet in report['sheets']:
        if len(prs.slides) >= MAX_POWERPOINT_SLIDES - 1: break
        slide = _base_slide(prs, f'ملخص الورقة: {sheet["name"]}', len(prs.slides) + 1)
        sheet_cards = [('الصفوف', f'{sheet["rows"]:,}', BLUE), ('الأعمدة', sheet['columns'], SKY), ('اكتمال البيانات', f'{sheet["completion_rate"]}%', TEAL), ('درجة الجودة', f'{sheet["quality_score"]}/100', AMBER)]
        for index, (label, value, color) in enumerate(sheet_cards):
            _kpi_card(slide, label, value, Inches(.72 + index * 3.04), Inches(1.25), color, width=Inches(2.8))
        add_bullet_panel(slide, 'أهم المشكلات', [x['issue'] for x in sheet['issues'][:3]] or ['لا توجد مشكلات بارزة.'], Inches(.72), Inches(2.82), Inches(5.85), Inches(3.62), max_items=3, color=RED if sheet['issues'] else TEAL)
        add_bullet_panel(slide, 'أهم النتائج', [f'متوسط {x["column"]}: {x["mean"]}' for x in sheet['numeric_stats'][:3]] or [f'عدد الخلايا: {sheet["cells"]:,}', f'القيم الفارغة: {sheet["missing_cells"]:,}'], Inches(6.82), Inches(2.82), Inches(5.9), Inches(3.62), max_items=3, color=BLUE)

    if len(prs.slides) < MAX_POWERPOINT_SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = NAVY
        for index, color in enumerate((BLUE, SKY, TEAL)):
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1 + index * .55), Inches(5.15 - index * .45), Inches(2.3), Inches(2.3))
            circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
        _textbox(slide, 'شكرًا لاستخدام Analytix', Inches(1), Inches(1.75), Inches(11.3), Inches(1), 34, True, WHITE, PP_ALIGN.CENTER)
        _textbox(slide, 'تم إعداد هذا التقرير اعتمادًا على البيانات المرفوعة ونتائج التحليل الآلي.', Inches(1.7), Inches(3.0), Inches(9.95), Inches(.75), 19, False, SKY, PP_ALIGN.CENTER)
        _textbox(slide, f'تاريخ التقرير: {report["generated_at"].strftime("%Y-%m-%d")}', Inches(2), Inches(4.05), Inches(9.3), Inches(.55), 16, False, WHITE, PP_ALIGN.CENTER)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
